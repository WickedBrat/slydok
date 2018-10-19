from pptx import Presentation
from googleapiclient import http
from googleapiclient.discovery import build
from httplib2 import Http
from oauth2client import file, client, tools
from io import BytesIO, FileIO

prs = Presentation()
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[0]

title.text = "Hello, World!"
subtitle.text = "python-pptx was here!"
prs.save('test.pptx')